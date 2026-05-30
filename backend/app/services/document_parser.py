"""文档解析与切块。

支持 PDF / DOCX / PPTX / TXT / MD，按页/幻灯片保留 page_number、section_title，
再做语义友好的切块（保留标题、控制大小、带 overlap）。
"""

from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path

SUPPORTED_SUFFIXES = {".pdf", ".docx", ".txt", ".md", ".pptx"}


@dataclass
class ParsedChunk:
    text: str
    chunk_index: int
    page_number: int | None = None
    section_title: str | None = None
    topic: str | None = None


def parse_document_pages(path: Path) -> list[tuple[int, str]]:
    """返回 [(page_number, text), ...]，页码从 1 开始。

    txt/md/docx 无天然分页，退化为单页全文。
    """
    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_SUFFIXES:
        raise ValueError(f"不支持的文件类型: {suffix}")
    if suffix == ".pdf":
        return _parse_pdf_pages(path)
    if suffix == ".pptx":
        return _parse_pptx_pages(path)
    if suffix == ".docx":
        return [(1, _parse_docx(path))]
    # txt / md
    return [(1, path.read_text(encoding="utf-8", errors="ignore"))]


def chunk_pages(
    pages: list[tuple[int, str]],
    min_size: int = 600,
    max_size: int = 1000,
    overlap: int = 150,
) -> list[ParsedChunk]:
    """把分页文本切块，全局递增 chunk_index，保留页码与小节标题。"""
    chunks: list[ParsedChunk] = []
    index = 0
    for page_number, page_text in pages:
        for piece, section_title in _chunk_text_with_titles(
            page_text, min_size=min_size, max_size=max_size, overlap=overlap
        ):
            chunks.append(
                ParsedChunk(
                    text=piece,
                    chunk_index=index,
                    page_number=page_number,
                    section_title=section_title,
                    topic=section_title or _first_line(piece),
                )
            )
            index += 1
    return chunks


# --------------------------------------------------------------------------- #
# 文本清洗与切块
# --------------------------------------------------------------------------- #


def clean_text(text: str) -> str:
    text = text.replace("　", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


_HEADING_RE = re.compile(r"^(#{1,4}\s+.+|[一二三四五六七八九十]+、.+|\d+[.、]\s*.+)$")


def _chunk_text_with_titles(
    text: str, min_size: int, max_size: int, overlap: int
) -> list[tuple[str, str | None]]:
    cleaned = clean_text(text)
    if not cleaned:
        return []

    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", cleaned) if p.strip()]
    results: list[tuple[str, str | None]] = []
    current = ""
    current_title: str | None = None

    for paragraph in paragraphs:
        if _HEADING_RE.match(paragraph.splitlines()[0].strip()):
            # 遇到新标题：先收束当前块
            if len(current) >= min_size:
                results.append((current, current_title))
                current = ""
            current_title = paragraph.splitlines()[0].strip("# 0123456789.、 ")

        if len(current) + len(paragraph) + 2 <= max_size:
            current = f"{current}\n\n{paragraph}".strip()
            continue

        if current:
            results.append((current, current_title))
        if len(paragraph) <= max_size:
            current = paragraph
        else:
            for piece in _split_long(paragraph, max_size, overlap):
                results.append((piece, current_title))
            current = ""

    if current:
        results.append((current, current_title))
    return results


def _split_long(text: str, max_size: int, overlap: int) -> list[str]:
    pieces = []
    step = max(1, max_size - overlap)
    for start in range(0, len(text), step):
        piece = text[start : start + max_size].strip()
        if piece:
            pieces.append(piece)
    return pieces


def _first_line(text: str) -> str | None:
    line = text.strip().splitlines()[0].strip("# 0123456789.、 ") if text.strip() else ""
    return line if 2 <= len(line) <= 40 else None


# --------------------------------------------------------------------------- #
# 各格式解析
# --------------------------------------------------------------------------- #


def _parse_pdf_pages(path: Path) -> list[tuple[int, str]]:
    try:
        import fitz  # PyMuPDF
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("解析 PDF 需要 pymupdf。") from exc

    pages: list[tuple[int, str]] = []
    with fitz.open(str(path)) as doc:
        for number, page in enumerate(doc, start=1):
            text = (page.get_text() or "").strip()
            if text:
                pages.append((number, text))
    return pages or [(1, "")]


def _parse_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("解析 DOCX 需要 python-docx。") from exc

    document = Document(str(path))
    return "\n\n".join(p.text for p in document.paragraphs if p.text.strip())


def _parse_pptx_pages(path: Path) -> list[tuple[int, str]]:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("解析 PPTX 需要 python-pptx。") from exc

    deck = Presentation(str(path))
    pages: list[tuple[int, str]] = []
    for number, slide in enumerate(deck.slides, start=1):
        texts = [
            shape.text
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if texts:
            pages.append((number, "\n\n".join(texts)))
    return pages or [(1, "")]
